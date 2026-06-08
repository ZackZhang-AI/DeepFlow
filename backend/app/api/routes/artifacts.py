"""
Artifact 生成 API — 播客脚本、PPT、文本操作
"""

import logging
import uuid
import re
import asyncio
import json
from abc import ABC, abstractmethod
from pathlib import Path
from urllib.parse import quote
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse, PlainTextResponse
from pydantic import BaseModel, Field

from backend.app.core.db import get_artifact, get_task, list_artifacts, save_artifact, save_report_version, update_task
from backend.app.core.auth import require_login
from cli.config import Config

router = APIRouter(prefix="/api/artifacts", tags=["artifacts"])

logger = logging.getLogger("deepflow.artifacts")
ARTIFACT_DIR = Config.OUTPUT_DIR / "artifacts"
ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)


class ArtifactRequest(BaseModel):
    task_id: str = Field(..., description="研究任务 ID")
    locale: str = Field(default="zh-CN")
    style: str = Field(default="general")  # for report re-generation


class ProsRequest(BaseModel):
    text: str = Field(..., min_length=10, description="待处理的文本")
    instruction: str = Field(default="")


class TTSProvider(ABC):
    """Audio synthesis provider contract for podcast artifacts."""

    name: str

    @abstractmethod
    def synthesize(self, text: str, output_path: Path) -> None:
        """Write synthesized audio to output_path."""


class LocalPyttsx3TTSProvider(TTSProvider):
    name = "local_pyttsx3"

    def synthesize(self, text: str, output_path: Path) -> None:
        import pyttsx3

        cleaned = _clean_tts_text(text)
        if not cleaned:
            raise RuntimeError("播客脚本文本为空，无法生成音频")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        engine = pyttsx3.init()
        engine.setProperty("rate", 165)
        engine.setProperty("volume", 0.95)
        engine.save_to_file(cleaned[:12000], str(output_path))
        engine.runAndWait()
        if not output_path.exists() or output_path.stat().st_size == 0:
            raise RuntimeError("本机 TTS 未生成音频文件")


class CloudTTSProvider(TTSProvider):
    """Reserved adapter shape for future cloud TTS implementations."""

    name = "cloud"

    def synthesize(self, text: str, output_path: Path) -> None:
        raise NotImplementedError("云 TTS Provider 尚未配置")


def get_tts_provider() -> TTSProvider:
    return LocalPyttsx3TTSProvider()


# ============================================================
# 播客脚本
# ============================================================

@router.post("/podcast")
async def generate_podcast(req: ArtifactRequest, user: dict = Depends(require_login)):
    """将报告转化为播客脚本"""
    task = get_task(req.task_id, user_id=user["user_id"])
    if not task or not task.get("report_markdown"):
        raise HTTPException(status_code=404, detail="报告不存在")

    import sys
    from pathlib import Path
    ROOT_DIR = Path(__file__).resolve().parent.parent.parent.parent.parent
    sys.path.insert(0, str(ROOT_DIR))

    from cli.agents.artifacts.podcast import generate_podcast_script, format_script_for_display

    title = task.get("plan_json", "")
    try:
        import json
        plan = json.loads(title) if title else {}
        title = plan.get("title", task["topic"])
    except Exception:
        title = task["topic"]

    script, pt, ct = await generate_podcast_script(
        report_markdown=task["report_markdown"],
        report_title=title,
        locale=req.locale,
    )

    if script is None:
        raise HTTPException(status_code=500, detail="播客脚本生成失败")

    display = format_script_for_display(script)
    audio_error = ""
    audio_artifact_id = None
    audio_url = None
    tts_provider = get_tts_provider()
    try:
        audio_path = ARTIFACT_DIR / f"{req.task_id}_podcast_{uuid.uuid4().hex[:8]}.wav"
        await asyncio.to_thread(tts_provider.synthesize, display, audio_path)
        audio_artifact = save_artifact(
            artifact_id=f"art_{uuid.uuid4().hex[:12]}",
            task_id=req.task_id,
            artifact_type="podcast_audio",
            title=f"{script.title}.wav",
            content=str(audio_path),
            metadata={"locale": req.locale, "format": "wav", "tts_provider": tts_provider.name},
            user_id=user["user_id"],
        )
        audio_artifact_id = audio_artifact["artifact_id"]
        audio_url = f"/api/artifacts/download/{audio_artifact_id}"
    except Exception as exc:
        audio_error = f"本机 TTS 音频生成失败，脚本已生成并可下载。原因：{exc}"
        logger.warning("Podcast audio generation failed: %s", exc)

    script_artifact = save_artifact(
        artifact_id=f"art_{uuid.uuid4().hex[:12]}",
        task_id=req.task_id,
        artifact_type="podcast",
        title=script.title,
        content=display,
        metadata={
            "tokens": pt + ct,
            "locale": req.locale,
            "audio_artifact_id": audio_artifact_id,
            "audio_error": audio_error,
            "tts_provider": tts_provider.name,
        },
        user_id=user["user_id"],
    )

    return {
        "artifact_id": script_artifact["artifact_id"],
        "audio_artifact_id": audio_artifact_id,
        "audio_url": audio_url,
        "audio_error": audio_error,
        "script": script.model_dump(),
        "display": display,
        "tokens": pt + ct,
    }


# ============================================================
# PPT
# ============================================================

@router.post("/ppt")
async def generate_ppt(req: ArtifactRequest, user: dict = Depends(require_login)):
    """将报告转化为 PPT Markdown"""
    task = get_task(req.task_id, user_id=user["user_id"])
    if not task or not task.get("report_markdown"):
        raise HTTPException(status_code=404, detail="报告不存在")

    import sys
    from pathlib import Path
    ROOT_DIR = Path(__file__).resolve().parent.parent.parent.parent.parent
    sys.path.insert(0, str(ROOT_DIR))

    from cli.agents.artifacts.ppt import compose_ppt

    title = task["topic"]
    try:
        import json
        plan = json.loads(task.get("plan_json", "{}"))
        title = plan.get("title", task["topic"])
    except Exception:
        pass

    slides, pt, ct = await compose_ppt(
        report_markdown=task["report_markdown"],
        report_title=title,
        locale=req.locale,
    )

    pptx_path = ARTIFACT_DIR / f"{req.task_id}_slides_{uuid.uuid4().hex[:8]}.pptx"
    _write_pptx(slides, title, pptx_path)
    pptx_artifact = save_artifact(
        artifact_id=f"art_{uuid.uuid4().hex[:12]}",
        task_id=req.task_id,
        artifact_type="pptx",
        title=f"{title}.pptx",
        content=str(pptx_path),
        metadata={"locale": req.locale, "format": "pptx"},
        user_id=user["user_id"],
    )
    markdown_artifact = save_artifact(
        artifact_id=f"art_{uuid.uuid4().hex[:12]}",
        task_id=req.task_id,
        artifact_type="ppt",
        title=title,
        content=slides,
        metadata={
            "tokens": pt + ct,
            "locale": req.locale,
            "slides_count": slides.count("---") + 1,
            "pptx_artifact_id": pptx_artifact["artifact_id"],
        },
        user_id=user["user_id"],
    )

    return {
        "artifact_id": markdown_artifact["artifact_id"],
        "pptx_artifact_id": pptx_artifact["artifact_id"],
        "pptx_url": f"/api/artifacts/download/{pptx_artifact['artifact_id']}",
        "slides_markdown": slides,
        "slides_count": slides.count("---") + 1,
        "tokens": pt + ct,
    }


# ============================================================
# 文本润色
# ============================================================

@router.post("/prose/improve")
async def improve_prose(req: ProsRequest, user: dict = Depends(require_login)):
    """润色文本"""
    import sys
    from pathlib import Path
    ROOT_DIR = Path(__file__).resolve().parent.parent.parent.parent.parent
    sys.path.insert(0, str(ROOT_DIR))

    from cli.agents.artifacts.prose import improve_text

    result, pt, ct = await improve_text(req.text, req.instruction)
    return {"result": result, "tokens": pt + ct}


@router.post("/prose/expand")
async def expand_prose(req: ProsRequest, user: dict = Depends(require_login)):
    """扩展文本"""
    import sys
    from pathlib import Path
    ROOT_DIR = Path(__file__).resolve().parent.parent.parent.parent.parent
    sys.path.insert(0, str(ROOT_DIR))

    from cli.agents.artifacts.prose import expand_text

    result, pt, ct = await expand_text(req.text)
    return {"result": result, "tokens": pt + ct}


@router.post("/prose/shorten")
async def shorten_prose(req: ProsRequest, user: dict = Depends(require_login)):
    """精简文本"""
    import sys
    from pathlib import Path
    ROOT_DIR = Path(__file__).resolve().parent.parent.parent.parent.parent
    sys.path.insert(0, str(ROOT_DIR))

    from cli.agents.artifacts.prose import shorten_text

    result, pt, ct = await shorten_text(req.text)
    return {"result": result, "tokens": pt + ct}


# ============================================================
# 多风格重生成
# ============================================================

@router.post("/restyle")
async def restyle_report(req: ArtifactRequest, user: dict = Depends(require_login)):
    """用指定风格重新生成报告"""
    task = get_task(req.task_id, user_id=user["user_id"])
    if not task:
        raise HTTPException(status_code=404, detail="任务不存在")

    import sys
    from pathlib import Path
    ROOT_DIR = Path(__file__).resolve().parent.parent.parent.parent.parent
    sys.path.insert(0, str(ROOT_DIR))

    from cli.agents.reporter import generate_report
    from cli.models import ResearchPlan, ResearchFinding

    # 重建 plan
    plan = ResearchPlan(title=task["topic"], locale=req.locale, has_enough_context=True, thought="", steps=[])
    try:
        import json
        pd = json.loads(task.get("plan_json", "{}"))
        plan = ResearchPlan(**pd)
    except Exception:
        pass

    # 简化版：直接复用缓存报告内容作为 findings 输入
    findings: list = []
    if task.get("report_markdown"):
        findings = [
            ResearchFinding(
                step_id="restyle_1",
                step_title="研究内容",
                problem_statement=task["topic"],
                findings_markdown=task["report_markdown"],
                conclusion="",
                references=[],
            )
        ]

    report, pt, ct = await generate_report(
        plan=plan,
        findings=findings,
        locale=req.locale,
        report_style=req.style,
    )

    artifact = save_artifact(
        artifact_id=f"art_{uuid.uuid4().hex[:12]}",
        task_id=req.task_id,
        artifact_type=f"report_style:{req.style}",
        title=f"{task['topic']} - {req.style}",
        content=report,
        metadata={"tokens": pt + ct, "locale": req.locale, "style": req.style},
        user_id=user["user_id"],
    )
    if task.get("report_markdown"):
        save_report_version(
            req.task_id,
            task["report_markdown"],
            user_id=user["user_id"],
            change_note=f"切换报告风格前版本: {req.style}",
        )
    update_task(
        req.task_id,
        owner_user_id=user["user_id"],
        report_markdown=report,
        tokens_used=(task.get("tokens_used") or 0) + pt + ct,
    )

    return {
        "artifact_id": artifact["artifact_id"],
        "style": req.style,
        "report_markdown": report,
        "tokens": pt + ct,
    }


@router.get("/download/{artifact_id}")
async def download_artifact(artifact_id: str, user: dict = Depends(require_login)):
    """下载已生成的成果物文件或文本内容。"""
    artifact = get_artifact(artifact_id, user_id=user["user_id"])
    if not artifact:
        raise HTTPException(status_code=404, detail="成果物不存在")

    path = _existing_file_path(artifact["content"])
    if path is None:
        if _is_text_artifact(artifact):
            filename = _artifact_filename(artifact)
            return PlainTextResponse(
                content=artifact["content"],
                media_type="text/markdown; charset=utf-8",
                headers={"Content-Disposition": _attachment_header(filename)},
            )
        raise HTTPException(status_code=404, detail="成果物文件不存在")

    suffix = path.suffix.lower()
    media_type = {
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".wav": "audio/wav",
    }.get(suffix, "application/octet-stream")
    return FileResponse(path, media_type=media_type, filename=path.name)


@router.get("/detail/{artifact_id}")
async def artifact_detail(artifact_id: str, user: dict = Depends(require_login)):
    """查看成果物详情；文本类成果物会返回正文。"""
    artifact = get_artifact(artifact_id, user_id=user["user_id"])
    if not artifact:
        raise HTTPException(status_code=404, detail="成果物不存在")
    return _serialize_artifact(artifact, include_content=True)


@router.get("/{task_id}")
async def list_task_artifacts(task_id: str, user: dict = Depends(require_login)):
    """列出某个任务已生成的成果物"""
    if not get_task(task_id, user_id=user["user_id"]):
        raise HTTPException(status_code=404, detail="任务不存在")
    return [_serialize_artifact(artifact) for artifact in list_artifacts(task_id, user_id=user["user_id"])]


def _write_pptx(slides_markdown: str, title: str, output_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    raw_slides = [s.strip() for s in re.split(r"\n---+\n", slides_markdown) if s.strip()]
    if not raw_slides:
        raw_slides = [f"# {title}\n\n暂无内容"]

    for index, raw in enumerate(raw_slides):
        lines = [line.strip() for line in raw.splitlines() if line.strip()]
        slide_title = title if index == 0 else f"Slide {index + 1}"
        body_lines: list[str] = []
        for line in lines:
            cleaned = re.sub(r"^#+\s*", "", line).strip()
            if line.startswith("#") and slide_title in (title, f"Slide {index + 1}"):
                slide_title = cleaned[:80] or slide_title
            else:
                body_lines.append(cleaned)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.7), Inches(5.5))
        body = body_box.text_frame
        body.word_wrap = True
        body.clear()
        for i, line in enumerate(body_lines[:12] or ["暂无正文"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = re.sub(r"^[-*]\s*", "", line)[:240]
            p.level = 1 if line.startswith(("-", "*")) else 0
            p.font.size = Pt(17 if len(p.text) < 80 else 14)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def _clean_tts_text(text: str) -> str:
    cleaned = re.sub(r"[*#_`>\[\]()]|https?://\S+", "", text)
    return re.sub(r"\s+", " ", cleaned).strip()


def _parse_metadata(artifact: dict) -> dict:
    raw = artifact.get("metadata_json") or "{}"
    if isinstance(raw, dict):
        return raw
    try:
        return json.loads(raw)
    except Exception:
        return {}


def _existing_file_path(content: str) -> Path | None:
    if not content or "\n" in content or len(content) > 500:
        return None
    try:
        path = Path(content)
        return path if path.exists() and path.is_file() else None
    except (OSError, ValueError):
        return None


def _is_text_artifact(artifact: dict) -> bool:
    artifact_type = artifact.get("artifact_type", "")
    return artifact_type in {"podcast", "ppt"} or artifact_type.startswith("report_style:")


def _artifact_filename(artifact: dict) -> str:
    title = re.sub(r'[\\/:*?"<>|\r\n]+', "_", artifact.get("title") or artifact["artifact_id"]).strip("._ ")
    suffix = ".md" if artifact.get("artifact_type") in {"podcast", "ppt"} else ".txt"
    if title.lower().endswith((".md", ".txt")):
        return title
    return f"{title or artifact['artifact_id']}{suffix}"


def _attachment_header(filename: str) -> str:
    ascii_name = filename.encode("ascii", "ignore").decode("ascii") or "artifact"
    return f'attachment; filename="{ascii_name}"; filename*=UTF-8\'\'{quote(filename)}'


def _serialize_artifact(artifact: dict, include_content: bool = False) -> dict:
    metadata = _parse_metadata(artifact)
    path = _existing_file_path(artifact.get("content", ""))
    text_artifact = _is_text_artifact(artifact)
    serialized = {
        **{k: v for k, v in artifact.items() if k != "content"},
        "metadata": metadata,
        "download_url": f"/api/artifacts/download/{artifact['artifact_id']}",
        "detail_url": f"/api/artifacts/detail/{artifact['artifact_id']}",
        "can_view": text_artifact,
        "is_file": path is not None,
    }
    if include_content:
        serialized["content"] = artifact["content"] if text_artifact else None
    return serialized
